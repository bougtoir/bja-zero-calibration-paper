#!/usr/bin/env python3
"""
Create BJA Figures & Tables PowerPoint — Japanese version
One figure or table per slide.
- Code-generated figures (1,2,4,5,6,7,8): embedded as PNG images
- Figure 3 (system comparison flow diagram): editable PowerPoint shapes
- Tables 1 & 2: editable PowerPoint tables
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGDIR = '/home/ubuntu/bja_figures'
OUTPATH = '/home/ubuntu/BJA_ZeroFree_Figures_JA.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_title_textbox(slide, text, top=Inches(0.2), font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return txBox


def add_caption_textbox(slide, text, top=Inches(6.6)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def add_image_slide(title, img_filename, caption, img_top=Inches(0.9)):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_title_textbox(slide, title)

    img_path = os.path.join(FIGDIR, img_filename)
    from PIL import Image
    im = Image.open(img_path)
    w_px, h_px = im.size
    aspect = w_px / h_px

    max_w = Inches(12.0)
    max_h = Inches(5.5)

    if aspect > (max_w / max_h):
        img_w = max_w
        img_h = int(max_w / aspect)
    else:
        img_h = max_h
        img_w = int(max_h * aspect)

    left = int((prs.slide_width - img_w) / 2)
    slide.shapes.add_picture(img_path, left, img_top, img_w, img_h)

    add_caption_textbox(slide, caption)
    return slide


def add_box(slide, left, top, width, height, fill_rgb, border_rgb=None, text='',
            font_size=Pt(11), font_color=RGBColor(0,0,0), bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for p in tf.paragraphs:
        tf._txBody.remove(p._p)
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.name = 'Arial'
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    tf.paragraphs[0].alignment = alignment

    shape.text_frame.margin_left = Pt(4)
    shape.text_frame.margin_right = Pt(4)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)

    return shape


def add_line_arrow(slide, x1, y1, x2, y2, color=RGBColor(0x33,0x33,0x33), width=Pt(2)):
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector._element.find(qn('a:ln'))
    if ln is None:
        from lxml import etree
        ln = etree.SubElement(connector._element.find(qn('p:spPr')), qn('a:ln'))
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return connector


def add_table_slide(slide, title_text, headers, data, caption_text,
                    col_widths=None, font_size=Pt(10), header_fill=RGBColor(0x2C,0x5F,0x8A)):
    add_title_textbox(slide, title_text)

    rows = len(data) + 1
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(12.3)
    height = Inches(0.4) * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                p.alignment = PP_ALIGN.CENTER
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_caption_textbox(slide, caption_text, top=Inches(1.2) + height + Inches(0.3))


# ══════════════════════════════════════════════════════════════
# スライド 1: 表紙
# ══════════════════════════════════════════════════════════════
slide0 = prs.slides.add_slide(BLANK_LAYOUT)
txBox = slide0.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '脈圧の正確性が暗示するゲインの妥当性：\nゼロ校正不要の観血的動脈圧モニタリングに向けた理論的枠組み'
p.font.size = Pt(24)
p.font.bold = True
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = '\nPulse Pressure Accuracy as an Implicit Gain Validator'
p2.font.size = Pt(14)
p2.font.italic = True
p2.font.name = 'Arial'
p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = '\n図表一覧'
p3.font.size = Pt(16)
p3.font.name = 'Arial'
p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p3.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════
# スライド 2: 図1 — 信号分解（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図1. 動脈圧波形の分解',
    'figure1_signal_decomposition.png',
    '図1. (A) 真の動脈圧：DC成分（オフセット依存のベースライン）とAC成分（脈圧、PP = 40 mmHg）。'
    '(B) DCオフセット（+15 mmHg）：PPは差分量のため40 mmHgで不変。'
    '(C) ゲインエラー（v = 1.15）：PPが46 mmHgに歪む。ゼロ校正は(B)を補正するが(C)は補正不可。'
)

# ══════════════════════════════════════════════════════════════
# スライド 3: 図2 — CCC コンコーダンスプロット（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図2. コンコーダンスプロット \u2014 4つのシミュレーションシナリオ',
    'figure2_ccc_zeroing_scenarios.png',
    '図2. (A) ゼロ校正前：オフセットによりCCC = 0.855。(B) ゼロ校正後：CCC = 0.986、Cb = 1.000。'
    '(C) ゲインエラー（v = 1.11）：Cb = 0.870、ゼロ校正では改善不可。(D) ゲイン＋オフセット：CCC = 0.976。'
    '破線 = 恒等線（y = x）；実線 = 回帰直線。'
)

# ══════════════════════════════════════════════════════════════
# スライド 4: 図3 — システム比較（編集可能シェイプ）
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK_LAYOUT)
add_title_textbox(slide3, '図3. 従来型 vs. ゼロ校正不要システムの比較')

# --- パネルA: 従来型システム（左半分） ---
txA = slide3.shapes.add_textbox(Inches(0.3), Inches(0.85), Inches(6), Inches(0.4))
tf = txA.text_frame
p = tf.paragraphs[0]
p.text = 'A. 従来の液充填システム'
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Arial'

# 動脈ボックス
add_box(slide3, Inches(0.5), Inches(1.6), Inches(1.6), Inches(0.9),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        '動脈\n（カテーテル先端）', Pt(11), bold=True)

# 矢印: 動脈 -> トランスデューサ
add_line_arrow(slide3, Inches(2.1), Inches(2.05), Inches(2.7), Inches(2.05),
               RGBColor(0x33, 0x66, 0xCC))

# ラベル: 液充填チューブ
tx_ft = slide3.shapes.add_textbox(Inches(1.9), Inches(1.35), Inches(1.2), Inches(0.3))
tf = tx_ft.text_frame
p = tf.paragraphs[0]
p.text = '液充填チューブ'
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0x33, 0x66, 0xCC)
p.font.name = 'Arial'

# 外部トランスデューサ
add_box(slide3, Inches(2.7), Inches(1.4), Inches(1.8), Inches(1.2),
        RGBColor(0xCC, 0xDD, 0xFF), RGBColor(0x33, 0x66, 0xCC),
        '外部\nトランスデューサ', Pt(12), bold=True)

# 矢印: トランスデューサ -> モニター
add_line_arrow(slide3, Inches(4.5), Inches(2.05), Inches(5.1), Inches(2.05))

# モニターボックス
add_box(slide3, Inches(5.1), Inches(1.5), Inches(1.3), Inches(1.0),
        RGBColor(0xDD, 0xDD, 0xDD), RGBColor(0x66, 0x66, 0x66),
        'モニター', Pt(12), bold=True)

# 3つのオフセット源ボックス
add_box(slide3, Inches(0.5), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        '静水圧カラム\n(\u0394P = \u03C1gh)', Pt(9))

add_box(slide3, Inches(2.2), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xDD, 0xCC), RGBColor(0xCC, 0x66, 0x00),
        '大気圧\n基準', Pt(9))

add_box(slide3, Inches(3.9), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xEE, 0xCC), RGBColor(0xCC, 0x99, 0x00),
        'トランスデューサ\nドリフト', Pt(9))

# 矢印 → 結果ボックス
add_line_arrow(slide3, Inches(1.25), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))
add_line_arrow(slide3, Inches(2.95), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))
add_line_arrow(slide3, Inches(4.65), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))

# 結果ボックス: 手動ゼロ校正必要
add_box(slide3, Inches(0.8), Inches(4.8), Inches(4.6), Inches(0.8),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        '手動ゼロ校正が必要\n（オフセット/uのみ補正；ゲイン/vは不変）',
        Pt(11), RGBColor(0xCC, 0x00, 0x00), bold=True)

# --- パネルB: 提案するゼロ校正不要システム（右半分） ---
txB = slide3.shapes.add_textbox(Inches(7.0), Inches(0.85), Inches(6), Inches(0.4))
tf = txB.text_frame
p = tf.paragraphs[0]
p.text = 'B. 提案するゼロ校正不要システム'
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(0x00, 0x66, 0x00)

# カテーテル先端MEMSボックス
add_box(slide3, Inches(7.2), Inches(1.4), Inches(2.0), Inches(1.2),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        'カテーテル先端\nMEMSセンサ\n（絶対圧）', Pt(11), bold=True)

# 矢印: MEMS -> スマートモニター
add_line_arrow(slide3, Inches(9.2), Inches(2.05), Inches(9.8), Inches(2.05),
               RGBColor(0x00, 0x88, 0x00))

# スマートモニターボックス
sm_shape = add_box(slide3, Inches(9.8), Inches(1.2), Inches(2.2), Inches(1.8),
        RGBColor(0xDD, 0xFF, 0xDD), RGBColor(0x00, 0x88, 0x00),
        '', Pt(11))

tf = sm_shape.text_frame
tf.clear()
p1 = tf.paragraphs[0]
p1.text = 'スマートモニター'
p1.font.size = Pt(12)
p1.font.bold = True
p1.font.name = 'Arial'
p2 = tf.add_paragraph()
p2.text = '\n内蔵気圧計\n＋自動補償'
p2.font.size = Pt(9)
p2.font.name = 'Arial'
p2.font.color.rgb = RGBColor(0x00, 0x66, 0x00)
p3 = tf.add_paragraph()
p3.text = '\n自己校正MEMS\n参照キャビティ'
p3.font.size = Pt(9)
p3.font.name = 'Arial'
p3.font.color.rgb = RGBColor(0x00, 0x66, 0x00)

# 矢印: スマートモニター -> ディスプレイ
add_line_arrow(slide3, Inches(12.0), Inches(2.1), Inches(12.3), Inches(2.1))

# ディスプレイボックス
add_box(slide3, Inches(12.3), Inches(1.5), Inches(0.8), Inches(1.0),
        RGBColor(0xDD, 0xDD, 0xDD), RGBColor(0x66, 0x66, 0x66),
        '表示', Pt(11), bold=True)

# 3つの排除ボックス
add_box(slide3, Inches(7.2), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        '先端センサ\n\u2192 \u0394h = 0\n\u2713 排除済み', Pt(9))

add_box(slide3, Inches(8.9), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xCC, 0xEE, 0xCC), RGBColor(0x00, 0x88, 0x00),
        '気圧補償\n\u2713 排除済み', Pt(9))

add_box(slide3, Inches(10.6), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xDD, 0xDD, 0xBB), RGBColor(0x00, 0x88, 0x00),
        '自己校正MEMS\n\u2713 排除済み', Pt(9))

# 矢印 → 結果ボックス
add_line_arrow(slide3, Inches(7.95), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))
add_line_arrow(slide3, Inches(9.65), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))
add_line_arrow(slide3, Inches(11.35), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))

# 結果ボックス: ゼロ校正不要
add_box(slide3, Inches(7.5), Inches(4.8), Inches(4.6), Inches(0.8),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        '設計によりゼロ校正不要\n（u = 0：設計で達成；PPの正確性がv = 1を検証）',
        Pt(11), RGBColor(0x00, 0x66, 0x00), bold=True)

# キャプション
add_caption_textbox(slide3,
    '図3. (A) 従来の液充填システム：3つのDCオフセット源が手動ゼロ校正を必要とし、'
    '校正はオフセット（u）のみを補正。(B) 提案するゼロ校正不要システム：カテーテル先端MEMSが静水圧カラムを排除、'
    '内蔵気圧計が大気圧を電子的に補償、自己校正MEMSがドリフトを排除。'
    '全オフセット源が設計により排除、PPの正確性がゲインを検証（v = 1）。',
    top=Inches(5.9))

# ══════════════════════════════════════════════════════════════
# スライド 5: 図4 — Cb診断空間（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図4. Cb診断空間（u\u2013v平面）',
    'figure4_cb_diagnostic_space.png',
    '図4. ロケーションシフト（u）とスケールシフト（v）の関数としてのバイアス補正係数（Cb）。'
    'ゼロ校正は機器を水平方向に移動（u \u2192 0）させるがvは変化しない。'
    'u = 0かつv = 1の場合にのみCb = 1.0を達成。点A\u2013Dは図2のシナリオに対応。'
)

# ══════════════════════════════════════════════════════════════
# スライド 6: 図5 — BA比較（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図5. Bland-Altmanプロット \u2014 4シナリオ',
    'figure5_ba_comparison.png',
    '図5. (A) ゼロ校正前：バイアス = 11.9 mmHg。(B) ゼロ校正後：バイアス \u2248 0。'
    '(C) ゲインエラー：バイアス = \u221210.9 mmHg、比例パターン。'
    '(D) ゲイン＋オフセット：バイアス \u2248 1.1 mmHg \u2014 BA上は良好に見えるがゲインエラーを隠蔽。'
    '赤色実線 = 平均バイアス；赤色破線 = 95%一致の限界。'
)

# ══════════════════════════════════════════════════════════════
# スライド 7: 図6 — 感度分析（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図6. 感度分析',
    'figure6_sensitivity_analysis.png',
    '図6. (A) ゲインエラー（%）とDCオフセット（mmHg）の関数としてのCb値ヒートマップ。'
    'ゼロ校正（オフセット = 0への移動）はゲインが正しい場合にのみCbを最大化。'
    '(B) 異なるセンサ精度（r = 0.95, 0.97, 0.99）でのゲインエラーに伴うCCC劣化。'
    '灰色帯域 = 典型的なMEMSゲイン精度範囲（\u00B15%）。'
)

# ══════════════════════════════════════════════════════════════
# スライド 8: 図7 — PP検証デモ（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図7. 脈圧（PP）検証のデモンストレーション',
    'figure7_pp_validation.png',
    '図7. (A) 正しいゲイン：基準波形（青）と測定波形（赤）が一致、PP = 40 mmHg。'
    '(B) ゲインエラー（v = 1.15）：PPが46 mmHgに拡大。(C) DCオフセットのみ（+15 mmHg）：PP不変。'
    '(D) 正しいゲインでのPP相関（傾き \u2248 1.0）。(E) ゲインエラーでのPP相関（傾き = 1.15）。'
    '(F) PPの正確性からゼロ校正不要モニタリングへの論理連鎖。'
)

# ══════════════════════════════════════════════════════════════
# スライド 9: 図8 — BA vs コンコーダンス（画像）
# ══════════════════════════════════════════════════════════════
add_image_slide(
    '図8. コンコーダンスプロット vs. Bland-Altmanプロット \u2014 並列比較',
    'figure8_ba_vs_concordance.png',
    '図8. 上段：コンコーダンスプロット。下段：Bland-Altmanプロット。同一データを2手法で可視化。'
    'シナリオBとDはBA上で類似のバイアス分布を示すが、コンコーダンスプロットでは'
    '回帰直線の傾きの違いとしてゲインエラーが明確に可視化される。'
    'CCC分解の補完的価値を実証。'
)

# ══════════════════════════════════════════════════════════════
# スライド 10: 表1 — DCオフセット源
# ══════════════════════════════════════════════════════════════
slide_t1 = prs.slides.add_slide(BLANK_LAYOUT)
add_table_slide(slide_t1,
    '表1. DCオフセットの源と工学的解決策',
    ['オフセット源', '物理的機序', '大きさ', '工学的解決策', 'CCC成分への影響'],
    [
        ['静水圧カラム', '\u0394P = \u03C1gh\n（トランスデューサ\u2013カテーテル先端高低差）', '~0.74 mmHg/cm',
         'カテーテル先端MEMSセンサ\n（液体カラムを排除）', 'uを低減\n（ロケーションシフト）'],
        ['大気圧基準', '大気（~760 mmHg）に\n対するゲージ圧測定', 'ベースライン全体',
         '絶対圧センサ＋\n内蔵気圧計', 'uを低減\n（ロケーションシフト）'],
        ['トランスデューサ\nドリフト', 'ストレインゲージの\n機械的クリープ・熱効果', '~1\u20135 mmHg/日',
         '内部参照圧力キャビティ付き\n自己校正MEMS', 'uを低減\n（ロケーションシフト）'],
        ['ゲインエラー\n（ゼロ校正で\n補正不可）', '感度不一致：\n出力/圧力 \u2260 公称値',
         '1\u201315%（典型的）', '工場校正；\nPPの正確性でゲイン検証',
         'vに影響\n（スケールシフト）；\nゼロ校正は無効'],
    ],
    '表1. 従来の観血的動脈圧モニタリングにおける3つのDCオフセット源と排除のための工学的解決策。'
    'ゲインエラー（最下行）はゼロ校正では補正不可であり、脈圧の正確性による別途の検証が必要。',
    font_size=Pt(9)
)

# ══════════════════════════════════════════════════════════════
# スライド 11: 表2 — 統計指標比較
# ══════════════════════════════════════════════════════════════
slide_t2 = prs.slides.add_slide(BLANK_LAYOUT)
add_table_slide(slide_t2,
    '表2. 4つのシミュレーションシナリオの統計指標比較（n = 150）',
    ['シナリオ', 'CCC', 'r', 'Cb', 'u', 'v', 'バイアス\n(mmHg)', 'LOA下限\n(mmHg)',
     'LOA上限\n(mmHg)', 'PE (%)'],
    [
        ['A: ゼロ校正前', '0.855', '0.986', '0.867', '\u22120.55', '0.99',
         '11.9', '4.8', '19.0', '7.2'],
        ['B: ゼロ校正後', '0.986', '0.986', '1.000', '0.01', '0.99',
         '\u22120.1', '\u22127.2', '7.0', '7.2'],
        ['C: ゲインエラー', '0.855', '0.982', '0.870', '0.54', '1.11',
         '\u221210.9', '\u221219.4', '\u22122.4', '8.7'],
        ['D: ゲイン＋\nオフセット', '0.976', '0.982', '0.993', '\u22120.05', '1.11',
         '1.1', '\u22127.4', '9.6', '8.7'],
    ],
    '表2. CCC = Linの一致性相関係数；r = ピアソンの相関係数（精度）；Cb = バイアス補正係数（正確度）；'
    'u = ロケーションシフト；v = スケールシフト；LOA = 一致の限界；PE = percentage error。'
    'シナリオBとCは類似したバイアス（\u22480）だがCbは大幅に異なり（1.000 vs. 0.870）、'
    'BA解析のみではゲインエラーを検出不可能であることを実証。',
    font_size=Pt(10)
)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
prs.save(OUTPATH)
print(f'Japanese PPTX saved: {OUTPATH}')
