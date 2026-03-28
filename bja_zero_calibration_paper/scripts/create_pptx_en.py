#!/usr/bin/env python3
"""
Create BJA Figures & Tables PowerPoint — English version
One figure or table per slide.
- Code-generated figures (1,2,4,5,6,7,8): embedded as PNG images
- Figure 3 (system comparison flow diagram): editable PowerPoint shapes
- Tables 1 & 2: editable PowerPoint tables
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGDIR = '/home/ubuntu/bja_figures'
OUTPATH = '/home/ubuntu/BJA_ZeroFree_Figures_EN.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_title_textbox(slide, text, top=Inches(0.2), font_size=Pt(18)):
    """Add a title text box at top of slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return txBox


def add_caption_textbox(slide, text, top=Inches(6.6)):
    """Add a caption text box at bottom of slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def add_image_slide(title, img_filename, caption, img_top=Inches(0.9)):
    """Add a slide with a centered image."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_title_textbox(slide, title)

    img_path = os.path.join(FIGDIR, img_filename)
    from PIL import Image
    im = Image.open(img_path)
    w_px, h_px = im.size
    aspect = w_px / h_px

    max_w = Inches(12.0)
    max_h = Inches(5.5)

    if aspect > (max_w / max_h):
        img_w = max_w
        img_h = int(max_w / aspect)
    else:
        img_h = max_h
        img_w = int(max_h * aspect)

    left = int((prs.slide_width - img_w) / 2)
    slide.shapes.add_picture(img_path, left, img_top, img_w, img_h)

    add_caption_textbox(slide, caption)
    return slide


def add_box(slide, left, top, width, height, fill_rgb, border_rgb=None, text='',
            font_size=Pt(11), font_color=RGBColor(0,0,0), bold=False, alignment=PP_ALIGN.CENTER):
    """Add a rounded rectangle shape with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for p in tf.paragraphs:
        tf._txBody.remove(p._p)
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.name = 'Arial'
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    tf.paragraphs[0].alignment = alignment

    shape.text_frame.margin_left = Pt(4)
    shape.text_frame.margin_right = Pt(4)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)

    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color_rgb=RGBColor(0x33,0x33,0x33)):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color_rgb
    connector.line.width = Pt(2)
    # Add arrowhead
    connector.begin_style = None
    connector.end_style = None
    line = connector.line
    line.fill.solid()
    line.fill.fore_color.rgb = color_rgb
    return connector


def add_line_arrow(slide, x1, y1, x2, y2, color=RGBColor(0x33,0x33,0x33), width=Pt(2)):
    """Add an arrow using a freeform or connector."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # Set end arrow
    ln = connector._element.find(qn('a:ln'))
    if ln is None:
        from lxml import etree
        ln = etree.SubElement(connector._element.find(qn('p:spPr')), qn('a:ln'))
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return connector


def add_table_slide(slide, title_text, headers, data, caption_text,
                    col_widths=None, font_size=Pt(10), header_fill=RGBColor(0x2C,0x5F,0x8A)):
    """Add an editable table to a slide."""
    add_title_textbox(slide, title_text)

    rows = len(data) + 1
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(12.3)
    height = Inches(0.4) * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill

    # Data rows
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                p.alignment = PP_ALIGN.CENTER
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_caption_textbox(slide, caption_text, top=Inches(1.2) + height + Inches(0.3))


# ══════════════════════════════════════════════════════════════
# Slide 1: Title slide
# ══════════════════════════════════════════════════════════════
slide0 = prs.slides.add_slide(BLANK_LAYOUT)
txBox = slide0.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'Pulse Pressure Accuracy as an Implicit Gain Validator:\nA Theoretical Framework for Zero-Calibration-Free Arterial Pressure Monitoring'
p.font.size = Pt(24)
p.font.bold = True
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = '\nFigures and Tables'
p2.font.size = Pt(16)
p2.font.name = 'Arial'
p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p2.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════
# Slide 2: Figure 1 — Signal Decomposition (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 1. Arterial Pressure Waveform Decomposition',
    'figure1_signal_decomposition.png',
    'Figure 1. (A) True arterial pressure with DC baseline and AC pulsatile component (PP = 40 mmHg). '
    '(B) DC offset (+15 mmHg): PP unchanged at 40 mmHg. (C) Gain error (v = 1.15): PP distorted to 46 mmHg. '
    'Zero calibration corrects (B) but not (C).'
)

# ══════════════════════════════════════════════════════════════
# Slide 3: Figure 2 — CCC Concordance Plots (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 2. Concordance Plots — Four Simulation Scenarios',
    'figure2_ccc_zeroing_scenarios.png',
    'Figure 2. (A) Before zeroing: offset causes CCC = 0.855. (B) After zeroing: CCC = 0.986, Cb = 1.000. '
    '(C) Gain error (v = 1.11): Cb = 0.870, uncorrectable by zeroing. (D) Gain + offset: CCC = 0.976. '
    'Dashed = identity line; solid = regression.'
)

# ══════════════════════════════════════════════════════════════
# Slide 4: Figure 3 — System Comparison (EDITABLE SHAPES)
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK_LAYOUT)
add_title_textbox(slide3, 'Figure 3. Conventional vs. Proposed Zero-Calibration-Free System')

# --- Panel A: Conventional system (left half) ---
panel_a_left = Inches(0.3)

# Panel A title
txA = slide3.shapes.add_textbox(panel_a_left, Inches(0.85), Inches(6), Inches(0.4))
tf = txA.text_frame
p = tf.paragraphs[0]
p.text = 'A. Conventional fluid-filled system'
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Arial'

# Artery box
add_box(slide3, Inches(0.5), Inches(1.6), Inches(1.6), Inches(0.9),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        'Artery\n(catheter tip)', Pt(11), bold=True)

# Arrow: Artery -> Transducer
add_line_arrow(slide3, Inches(2.1), Inches(2.05), Inches(2.7), Inches(2.05),
               RGBColor(0x33, 0x66, 0xCC))

# Label: Fluid-filled tubing
tx_ft = slide3.shapes.add_textbox(Inches(1.9), Inches(1.35), Inches(1.2), Inches(0.3))
tf = tx_ft.text_frame
p = tf.paragraphs[0]
p.text = 'Fluid-filled tubing'
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0x33, 0x66, 0xCC)
p.font.name = 'Arial'

# External transducer box
add_box(slide3, Inches(2.7), Inches(1.4), Inches(1.8), Inches(1.2),
        RGBColor(0xCC, 0xDD, 0xFF), RGBColor(0x33, 0x66, 0xCC),
        'External\ntransducer', Pt(12), bold=True)

# Arrow: Transducer -> Monitor
add_line_arrow(slide3, Inches(4.5), Inches(2.05), Inches(5.1), Inches(2.05))

# Monitor box
add_box(slide3, Inches(5.1), Inches(1.5), Inches(1.3), Inches(1.0),
        RGBColor(0xDD, 0xDD, 0xDD), RGBColor(0x66, 0x66, 0x66),
        'Monitor', Pt(12), bold=True)

# Three offset source boxes
add_box(slide3, Inches(0.5), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        'Hydrostatic\ncolumn\n(\u0394P = \u03C1gh)', Pt(9))

add_box(slide3, Inches(2.2), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xDD, 0xCC), RGBColor(0xCC, 0x66, 0x00),
        'Atmospheric\npressure\nreference', Pt(9))

add_box(slide3, Inches(3.9), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xFF, 0xEE, 0xCC), RGBColor(0xCC, 0x99, 0x00),
        'Transducer\ndrift', Pt(9))

# Arrows down to result box
add_line_arrow(slide3, Inches(1.25), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))
add_line_arrow(slide3, Inches(2.95), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))
add_line_arrow(slide3, Inches(4.65), Inches(4.1), Inches(2.95), Inches(4.8),
               RGBColor(0xCC, 0x00, 0x00))

# Result box: Manual zero calibration required
add_box(slide3, Inches(0.8), Inches(4.8), Inches(4.6), Inches(0.8),
        RGBColor(0xFF, 0xCC, 0xCC), RGBColor(0xCC, 0x00, 0x00),
        'Manual zero calibration required\n(corrects offset/u only; gain/v unchanged)',
        Pt(11), RGBColor(0xCC, 0x00, 0x00), bold=True)

# --- Panel B: Proposed zero-free system (right half) ---
panel_b_left = Inches(7.0)

# Panel B title
txB = slide3.shapes.add_textbox(panel_b_left, Inches(0.85), Inches(6), Inches(0.4))
tf = txB.text_frame
p = tf.paragraphs[0]
p.text = 'B. Proposed zero-free system'
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(0x00, 0x66, 0x00)

# Catheter-tip MEMS box
add_box(slide3, Inches(7.2), Inches(1.4), Inches(2.0), Inches(1.2),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        'Catheter-tip\nMEMS sensor\n(absolute pressure)', Pt(11), bold=True)

# Arrow: MEMS -> Smart Monitor
add_line_arrow(slide3, Inches(9.2), Inches(2.05), Inches(9.8), Inches(2.05),
               RGBColor(0x00, 0x88, 0x00))

# Smart Monitor box (larger, with sub-text)
sm_shape = add_box(slide3, Inches(9.8), Inches(1.2), Inches(2.2), Inches(1.8),
        RGBColor(0xDD, 0xFF, 0xDD), RGBColor(0x00, 0x88, 0x00),
        '', Pt(11))

# Smart Monitor internal text
tf = sm_shape.text_frame
tf.clear()
p1 = tf.paragraphs[0]
p1.text = 'Smart Monitor'
p1.font.size = Pt(12)
p1.font.bold = True
p1.font.name = 'Arial'
p2 = tf.add_paragraph()
p2.text = '\nBuilt-in barometer\n+ auto-compensation'
p2.font.size = Pt(9)
p2.font.name = 'Arial'
p2.font.color.rgb = RGBColor(0x00, 0x66, 0x00)
p3 = tf.add_paragraph()
p3.text = '\nSelf-calibrating\nMEMS reference'
p3.font.size = Pt(9)
p3.font.name = 'Arial'
p3.font.color.rgb = RGBColor(0x00, 0x66, 0x00)

# Arrow: Smart Monitor -> Display
add_line_arrow(slide3, Inches(12.0), Inches(2.1), Inches(12.3), Inches(2.1))

# Display box
add_box(slide3, Inches(12.3), Inches(1.5), Inches(0.8), Inches(1.0),
        RGBColor(0xDD, 0xDD, 0xDD), RGBColor(0x66, 0x66, 0x66),
        'Display', Pt(11), bold=True)

# Three elimination boxes
add_box(slide3, Inches(7.2), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        'Tip sensor\n\u2192 \u0394h = 0\n\u2713 Eliminated', Pt(9))

add_box(slide3, Inches(8.9), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xCC, 0xEE, 0xCC), RGBColor(0x00, 0x88, 0x00),
        'Barometric\ncompensation\n\u2713 Eliminated', Pt(9))

add_box(slide3, Inches(10.6), Inches(3.2), Inches(1.5), Inches(0.9),
        RGBColor(0xDD, 0xDD, 0xBB), RGBColor(0x00, 0x88, 0x00),
        'Self-calibrating\nMEMS\n\u2713 Eliminated', Pt(9))

# Arrows down to result box
add_line_arrow(slide3, Inches(7.95), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))
add_line_arrow(slide3, Inches(9.65), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))
add_line_arrow(slide3, Inches(11.35), Inches(4.1), Inches(9.65), Inches(4.8),
               RGBColor(0x00, 0x88, 0x00))

# Result box: Zero calibration unnecessary
add_box(slide3, Inches(7.5), Inches(4.8), Inches(4.6), Inches(0.8),
        RGBColor(0xCC, 0xFF, 0xCC), RGBColor(0x00, 0x88, 0x00),
        'Zero calibration unnecessary by design\n(u = 0 by design; PP accuracy validates v = 1)',
        Pt(11), RGBColor(0x00, 0x66, 0x00), bold=True)

# Caption
add_caption_textbox(slide3,
    'Figure 3. (A) Conventional fluid-filled system: three DC offset sources require manual zero calibration, '
    'which corrects offset (u) only. (B) Proposed zero-free system: catheter-tip MEMS eliminates hydrostatic column, '
    'built-in barometer compensates atmospheric pressure, self-calibrating MEMS eliminates drift. '
    'All offset sources eliminated by design; PP accuracy validates gain (v = 1).',
    top=Inches(5.9))


# ══════════════════════════════════════════════════════════════
# Slide 5: Figure 4 — Cb diagnostic space (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 4. Cb Diagnostic Space (u\u2013v Plane)',
    'figure4_cb_diagnostic_space.png',
    'Figure 4. Bias correction factor (Cb) as a function of location shift (u) and scale shift (v). '
    'Zero calibration moves the device horizontally (u \u2192 0) but does not change v. '
    'Only u = 0 and v = 1 achieves Cb = 1.0. Points A\u2013D correspond to Figure 2 scenarios.'
)

# ══════════════════════════════════════════════════════════════
# Slide 6: Figure 5 — BA comparison (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 5. Bland-Altman Plots \u2014 Four Scenarios',
    'figure5_ba_comparison.png',
    'Figure 5. (A) Before zeroing: bias = 11.9 mmHg. (B) After zeroing: bias \u2248 0. '
    '(C) Gain error: bias = \u221210.9 mmHg with proportional pattern. '
    '(D) Gain + offset: bias \u2248 1.1 mmHg \u2014 appears acceptable on BA but hides gain error. '
    'Red solid = mean bias; red dashed = 95% limits of agreement.'
)

# ══════════════════════════════════════════════════════════════
# Slide 7: Figure 6 — Sensitivity analysis (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 6. Sensitivity Analysis',
    'figure6_sensitivity_analysis.png',
    'Figure 6. (A) Cb heatmap as a function of gain error (%) and DC offset (mmHg). '
    'Zero calibration (horizontal move to offset = 0) only maximizes Cb when gain is correct. '
    '(B) CCC degradation with gain error at different sensor precision levels (r = 0.95, 0.97, 0.99). '
    'Grey band = typical MEMS gain accuracy range (\u00B15%).'
)

# ══════════════════════════════════════════════════════════════
# Slide 8: Figure 7 — PP validation demo (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 7. Pulse Pressure Validation Demonstration',
    'figure7_pp_validation.png',
    'Figure 7. (A) Correct gain: reference (blue) and measured (red) match, PP = 40 mmHg. '
    '(B) Gain error (v = 1.15): PP expanded to 46 mmHg. (C) DC offset only (+15 mmHg): PP unchanged. '
    '(D) Correct gain PP correlation (slope \u2248 1.0). (E) Gain error PP correlation (slope = 1.15). '
    '(F) Logical chain from PP accuracy to zero-calibration-free monitoring.'
)

# ══════════════════════════════════════════════════════════════
# Slide 9: Figure 8 — BA vs Concordance (image)
# ══════════════════════════════════════════════════════════════
add_image_slide(
    'Figure 8. Concordance Plots vs. Bland-Altman Plots \u2014 Side-by-Side',
    'figure8_ba_vs_concordance.png',
    'Figure 8. Top row: concordance plots. Bottom row: Bland-Altman plots. Same data, two methods. '
    'Scenarios B and D show similar BA bias distributions, but concordance plots reveal '
    'gain error as regression slope deviation. Demonstrates the complementary value of CCC decomposition.'
)

# ══════════════════════════════════════════════════════════════
# Slide 10: Table 1 — DC offset sources
# ══════════════════════════════════════════════════════════════
slide_t1 = prs.slides.add_slide(BLANK_LAYOUT)
add_table_slide(slide_t1,
    'Table 1. DC Offset Sources and Engineering Solutions',
    ['Offset Source', 'Physical Mechanism', 'Magnitude', 'Engineering Solution', 'CCC Component'],
    [
        ['Hydrostatic column', '\u0394P = \u03C1gh\n(transducer\u2013catheter tip height)', '~0.74 mmHg/cm',
         'Catheter-tip MEMS sensor\n(eliminates fluid column)', 'Reduces u\n(location shift)'],
        ['Atmospheric pressure\nreference', 'Gauge measurement\nrelative to ~760 mmHg', 'Entire baseline',
         'Absolute pressure sensor +\nbuilt-in barometer', 'Reduces u\n(location shift)'],
        ['Transducer drift', 'Mechanical creep,\nthermal effects', '~1\u20135 mmHg/day',
         'Self-calibrating MEMS\nwith reference cavity', 'Reduces u\n(location shift)'],
        ['Gain error\n(NOT corrected by\nzero calibration)', 'Sensitivity mismatch:\noutput/pressure \u2260 nominal',
         '1\u201315% (typical)', 'Factory calibration;\nPP accuracy validates gain',
         'Affects v\n(scale shift);\nzeroing ineffective'],
    ],
    'Table 1. Three DC offset sources in conventional invasive arterial pressure monitoring '
    'and engineering solutions for elimination. Note that gain error (bottom row) is NOT correctable '
    'by zero calibration and requires separate validation via pulse pressure accuracy.',
    font_size=Pt(9)
)

# ══════════════════════════════════════════════════════════════
# Slide 11: Table 2 — Statistical comparison
# ══════════════════════════════════════════════════════════════
slide_t2 = prs.slides.add_slide(BLANK_LAYOUT)
add_table_slide(slide_t2,
    'Table 2. Statistical Comparison Across Four Simulation Scenarios (n = 150)',
    ['Scenario', 'CCC', 'r', 'Cb', 'u', 'v', 'Bias\n(mmHg)', 'LOA low\n(mmHg)',
     'LOA high\n(mmHg)', 'PE (%)'],
    [
        ['A: Before zeroing', '0.855', '0.986', '0.867', '\u22120.55', '0.99',
         '11.9', '4.8', '19.0', '7.2'],
        ['B: After zeroing', '0.986', '0.986', '1.000', '0.01', '0.99',
         '\u22120.1', '\u22127.2', '7.0', '7.2'],
        ['C: Gain error', '0.855', '0.982', '0.870', '0.54', '1.11',
         '\u221210.9', '\u221219.4', '\u22122.4', '8.7'],
        ['D: Gain + offset', '0.976', '0.982', '0.993', '\u22120.05', '1.11',
         '1.1', '\u22127.4', '9.6', '8.7'],
    ],
    'Table 2. CCC = Lin\'s concordance correlation coefficient; r = Pearson correlation (precision); '
    'Cb = bias correction factor (accuracy); u = location shift; v = scale shift; '
    'LOA = limits of agreement; PE = percentage error. '
    'Scenarios B and C have similar bias (\u22480) but very different Cb (1.000 vs. 0.870), '
    'demonstrating that BA analysis alone cannot detect gain error.',
    font_size=Pt(10)
)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
prs.save(OUTPATH)
print(f'English PPTX saved: {OUTPATH}')
